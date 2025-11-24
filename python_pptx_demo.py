#%%Libraries
import pandas as pd
import numpy as np
import itertools
import matplotlib.pyplot as plt
import seaborn as sns
from io import BytesIO
from pptx import Presentation, util
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.chart.shared import CT_Layout
from pptx.chart.data import ChartData

#%% Config Variables
primary_color = {'r':7,'g':24,'b':45}               # Navy Blue
secondary_color = {'r':241,'g':221,'b':157}         # Gold
white_color = {'r':255,'g':255,'b':255}             # White
positive_color = {'r':204,'g':255,'b':204}          # Green
positive_accent_color = {'r':0,'g':108,'b':49}      # Dark Green
negative_color = {'r':255,'g':209,'b':209}          # Red
negative_accent_color = {'r':155,'g':17,'b':27}     # Dark Red
standard_column_width = 1.00
pivot_column_width = 1.30
standard_row_height = 0.30

#%% Functions

def SubElement(parent, tagname, **kwargs):
    '''
    Function to define an element (common in XML syntax)
    '''
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _horizontal_alignment(shape, horizontal_alignment:str):
    '''
    Text align a paragraph horizontally.
    horizontal_alignment: alignment string; default is typically 'left' (dependent upon default settings)
    '''
    match horizontal_alignment:
        case 'left': shape.alignment = PP_ALIGN.LEFT
        case 'center': shape.alignment = PP_ALIGN.CENTER
        case 'right': shape.alignment = PP_ALIGN.RIGHT

def _vertical_alignment(shape, vertical_alignment:str):
    '''
    Text align a textframe vertically.
    vertical_alignment: alignment string; default is typically 'top' (dependent upon default settings)
    '''
    match vertical_alignment:
        case 'top': shape.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        case 'middle': shape.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        case 'bottom': shape.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM

def _shape_fill(shape, fill_color:dict):
    '''
    Fill shape with a solid color.
    fill_color: dictionary containing desired 'r', 'g', and 'b' color values
    '''
    ##Access shape's fill property
    fill = shape.fill
    ##Set fill pattern to solid
    fill.solid()
    ##Set shape's fill color as RGB value
    fill.fore_color.rgb = RGBColor(fill_color['r'],fill_color['g'],fill_color['b'])

def _shape_color(shape, line_color:dict):
    '''
    Fill shape outline with a solid color.
    line_color: dictionary containing desired 'r', 'g', and 'b' color values
    '''
    ##Access shape's line property
    line = shape.line
    ##Set shape's line color as RGB value
    line.color.rgb = RGBColor(line_color['r'],line_color['g'],line_color['b'])

def _paragraphs(shape, text:list, font_size:int, font_color:dict, bold:bool, horizontal_alignment:str, vertical_alignment:str, bulleted:bool, run_bold:list):
    '''
    Creates paragraph object
    '''
    ##Access shape's text frame
    tf = shape.text_frame
    ##A text frame has an initial paragraph object; we want to delete that and add our own
    if tf.paragraphs[0]: tf.paragraphs[0]._element.getparent().remove(tf.paragraphs[0]._element)
    ##Ensure text will maintain set width
    tf.word_wrap = True
    ##Loop through provided list of text; each item will be treated as its own paragraph
    for line in text:
        ##Add paragraph to shape's text frame
        p = tf.add_paragraph()
        ##If the item is also a list, break sub-items into separate runs to apply different text formatting
        if isinstance(line, list):
            for index, sub_line in enumerate(line):
                ##Add run to shape's paragraph
                run = p.add_run()
                ##Insert text as string
                run.text = sub_line
                run.font.size = util.Pt(font_size)
                run.font.bold = run_bold[index]
                run.font.color.rgb = RGBColor(font_color['r'],font_color['g'],font_color['b'])
        else:
            run = p.add_run()
            run.text = line
            run.font.size = util.Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(font_color['r'],font_color['g'],font_color['b'])
        ##If bulleted=True, make each paragraph in text frame a bullet point
        if bulleted:
            ##Access paragraph properties
            pPr = p._p.get_or_add_pPr()
            ##Set indentation between text frame and bullet; number is in EMU
            pPr.set('marL','0')
            ##Set indentation between bullet and paragraph; number is in EMU
            pPr.set('indent','171450')
            ##Define bullet properties, including the character symbol
            _ = SubElement(parent=pPr, tagname='a:buChar', char='●')
        ##Set text alignment
        _horizontal_alignment(p, horizontal_alignment)
        _vertical_alignment(tf, vertical_alignment)

def _cell_margins(shape, margin_left:float=0, margin_right:float=0, margin_top:float=0, margin_bottom:float=0, vertical_alignment:str='middle'):
    '''
    Set table cell margins.
    margin_left, margin_right, margin_top, margin_bottom: inches of padding within a table cell
    vertical_alignment: alignment string; default is typically 'top' (dependent upon default settings)
    '''
    shape.margin_left = Inches(margin_left)
    shape.margin_right = Inches(margin_right)
    shape.margin_top = Inches(margin_top)
    shape.margin_bottom = Inches(margin_bottom)
    _vertical_alignment(shape, vertical_alignment)

def _cell_border(cell, borders:list, line_type:str='solid', border_color:str="07182D", border_width:str='12700'):
    '''
    Apply table cell borders.
    borders: desired cell borders being accessed ('a:lnL'= left, 'a:lnR'= right, 'a:lnT'= top, 'a:lnB'= bottom); 'a:lnB' is the only side that works consistently
    line_type: border style ('solid', 'dash', 'dot', 'lgDash', 'dashDot', 'lgDashDot', 'sysDash', 'sysDot', 'sysDashDot', 'sysDashDotDot')
    border_color: border color; any hex code value
    border_width: string of English Metric Units (EMU), a common size standardization; 1px = 12700 EMUs, 1in = 914400 EMUs
    '''
    fill_found = False
    fill = cell.fill
    ##Check if cell already has a fill color; to properly apply a cell border, the background must be transparent
    if cell.fill.type == MSO_FILL.SOLID:
        ##Store fill color to reapply later
        fill_color = fill.fore_color.rgb
        fill_found = True
        ##Set cell fill to transparent
        fill.background()
    ##Access cell attributes
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    ##Clean up existing cell formatting by removing default cell attributes
    for lines in borders:
        tag = lines.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag):
                tcPr.remove(e)
        ##Apply desired cell borders
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        ##Apply border color and style
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val=line_type)
        ##Additional formatting options
        # round_ = SubElement(ln, 'a:round')
        # headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        # tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
    ##Reapply fill color if any was present
    if fill_found:
        fill.solid()
        cell.fill.fore_color.rgb = fill_color
    return cell

def _table_styles(shape):
    '''
    Apply specific table styling.
    Desired style is 'No Style, No Grid'
    Additional style GUIDs can be found at: https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372
    '''
    table_style = shape._element.graphic.graphicData.tbl
    table_style[0][-1].text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'

def _table_column_sizes(table, widths):
    '''
    Set each column width to the corresponding list of sizes
    '''
    table_columns = table.columns
    for index, width in enumerate(widths):
        table_columns[index].width = Inches(width)

def _table_row_sizes(table, heights):
    '''
    Set each row height to the corresponding list of sizes
    '''
    table_rows = table.rows
    for index, height in enumerate(heights):
        table_rows[index].height = Inches(height)

def _plot_area_only(chart):
    '''
    Removes all additional plot elements.
    Used when creating sparklines
    '''
    ##Remove y-axis gridlines
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    ##Remove y-axis title
    chart.value_axis.has_title = False
    ##Hide y-axis
    chart.value_axis.visible = False
    ##Set y-labels to a small font size and color same as slide background
    chart.value_axis.tick_labels.font.size = util.Pt(1)
    chart.value_axis.tick_labels.font.color.rgb = RGBColor(255,255,255)
    ##Remove x-axis gridlines
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    ##Remove x-axis title
    chart.category_axis.has_title = False
    ##Hide x-axis
    chart.category_axis.visible = False
    ##Set x-labels to a small font size and color same as slide background
    chart.category_axis.tick_labels.font.size = util.Pt(1)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor(255,255,255)
    ##Remove x-labels
    chart.category_axis.has_labels = False
    ##Remove chart title
    chart.has_title = False
    ##Remove legend
    chart.has_legend = False
    ##Remove any other plot element
    plot_area = chart.element.chart.plotArea
    plot_area_children = plot_area.getchildren()
    pac_layout = [x for x in plot_area_children if isinstance(x, CT_Layout)]
    if len(pac_layout) == 1:
        layout = pac_layout[0]
    else:
        layout = SubElement(plot_area, 'c:layout', val='inner')
    manual_layout = layout.get_or_add_manualLayout()
    manual_layout.remove_all(*['c:layoutTarget','c:xMode','c:yMode','c:x','c:y','c:w','c:h'])
    _ = SubElement(manual_layout, 'c:layoutTarget', val='inner')
    _ = SubElement(manual_layout, 'c:xMode', val='edge')
    _ = SubElement(manual_layout, 'c:yMode', val='edge')
    _ = SubElement(manual_layout, 'c:x', val=str(0))
    _ = SubElement(manual_layout, 'c:y', val=str(0))
    _ = SubElement(manual_layout, 'c:w', val=str(1))
    _ = SubElement(manual_layout, 'c:h', val=str(1))

def autoshape(shapes, mso_shape, left:float, top:float, width:float, height:float, fill_color:dict=primary_color, line_color:dict=white_color, no_border:bool=False, line_width:float=util.Pt(1), rotate:float=0, shadow:bool=True):
    '''
    Create an autoshape object.
    mso_shape: desired autoshape option
    no_border: if true, the outline will be transparent
    line_width: border thickness
    rotate: number of degrees to rotate the autoshape
    shadow: whether or not to include the default shadow effect
    '''
    shape = shapes.add_shape(mso_shape, left, top, width, height)
    _shape_fill(shape, fill_color=fill_color)
    _shape_color(shape, line_color=line_color)
    if no_border: shape.line.fill.background()
    shape.shadow.inherit = shadow
    shape.line.width = line_width
    shape.rotation = rotate
    return shape

def shape_text(shape, text:list, font_size:int, bold:bool=False, font_color:dict=primary_color, fill_color=False, border_color=False, horizontal_alignment:str='left', vertical_alignment:str='top', bulleted:bool=False, run_bold:list=[False,False]):
    '''
    Add text to a shape object.
    bulleted: whether or not to make each line of text in list as a bulleted item
    run_bold: bold two different runs in a bulleted line; useful for bolding a bulleted topic, and unbold its definition
    '''
    _paragraphs(shape, text, font_size, font_color, bold, horizontal_alignment, vertical_alignment, bulleted, run_bold)
    if border_color: _shape_color(shape, border_color)
    if fill_color: _shape_fill(shape, fill_color)

def table(slide, df, title:str, horizontal_pos:float, vertical_pos:float, font_size:float=10):
    '''
    Creates table object
    '''
    ##Create table with (number of rows, number of columns, left position, top position, width, height)
    shape = slide.shapes.add_table(len(df)+2, len(df.columns), Inches(horizontal_pos), Inches(vertical_pos), Inches((standard_column_width*(len(df.columns)-1))+pivot_column_width), Inches(standard_row_height*(len(df)+2)))
    table = shape.table
    ##Apply table styling
    _table_styles(shape)
    ##Set column widths. First column is wider and the rest are standard column width
    _table_column_sizes(table, [pivot_column_width]+[standard_column_width]*(len(df.columns)-1))
    ##Access first cell in table (base 0)
    cell = table.cell(0,0)
    ##Access last cell in the first row
    other_cell = table.cell(0,len(df.columns)-1)
    ##Merge all cells between them into a single cell
    cell.merge(other_cell)
    ##Set cell margins
    _cell_margins(cell)
    ##Cell color background
    _shape_fill(cell, fill_color=primary_color)
    ##Insert text into cell
    shape_text(cell, [title], font_size, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
    ##Loop through each cell in the dataframe and insert into the table
    for i in range(1,len(df)+2):
        for j in range(0,len(df.columns)):
            cell = table.cell(i,j)
            _cell_margins(cell)
            ##If the cell is in the second row, insert column header
            if i == 1:
                _shape_fill(cell, fill_color=primary_color)
                shape_text(cell, [str(df.columns[j])], font_size, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
            ##If the dataframe cell is a number, apply certain formatting
            elif isinstance(df.iloc[i-2,j], (int, float)):
                cell.margin_right = Inches(0.05)
                shape_text(cell, [f'$ {df.iloc[i-2,j]:,.0f}'], font_size, bold=i-1==len(df), horizontal_alignment='right', vertical_alignment='middle')
                ##If cell is in the last or second to last row, apply a bottom cell border
                if (i-1 == len(df)) or (i-1 == len(df)-1): cell = _cell_border(cell, ['a:lnB'])
            ##If the dataframe cell is a string, apply generic formatting
            else:
                cell.margin_left = Inches(0.05)
                shape_text(cell, [str(df.iloc[i-2,j])], font_size, bold=i-1==len(df), horizontal_alignment='left', vertical_alignment='middle')
                if (i-1 == len(df)) or (i-1 == len(df)-1): cell = _cell_border(cell, ['a:lnB'])
    ##Set row heights. Must be set after text is inserted into all cells. When text is initially inserted into a cell it has the default font size; if the desired row height is smaller than default font size, the cell will not get smaller.
    _table_row_sizes(table, [standard_row_height]*(len(df)+2))

def trend_table(slide, df, first_line_title, second_line_title, horizontal_pos, vertical_pos, chart_type, font_size:float=10):
    '''
    Creates placeholder table where sparklines will be overlaid
    '''
    shape = slide.shapes.add_table(len(df)+2, 1, Inches(horizontal_pos), Inches(vertical_pos), Inches(pivot_column_width), Inches(standard_row_height*(len(df)+2)))
    table = shape.table
    _table_styles(shape)
    _table_column_sizes(table, [pivot_column_width])
    ##Table headers
    for index, title in enumerate([first_line_title, second_line_title]):
        cell = table.cell(index, 0)
        _cell_margins(cell)
        _shape_fill(cell, fill_color=primary_color)
        shape_text(cell, [title], font_size, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
    ##To maintain desired formatting, text must be in every cell; '0' is inserted into each cell and colored the same as the slide background to appear hidden.
    for i in range(2, len(df)+2):
        cell = table.cell(i,0)
        _cell_margins(cell)
        shape_text(cell, ['0'], font_size, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
        if (i==len(df)) or (i==len(df)+1): cell = _cell_border(cell, ['a:lnB'])
    _table_row_sizes(table, [standard_row_height]*(len(df)+2))
    ##Overlay sparkline/sparkbar in the table cell
    for index, row in df.iterrows():
        if chart_type == 'line': sparkline(slide, [item for item in row.tolist() if isinstance(item, (int, float))], Inches(horizontal_pos), Inches(vertical_pos+(standard_row_height*2)+(standard_row_height*index)), Inches(pivot_column_width), Inches(standard_row_height))
        if chart_type == 'bar': sparkbar(slide, [item for item in row.tolist() if isinstance(item, (int, float))], Inches(horizontal_pos), Inches(vertical_pos+(standard_row_height*2)+(standard_row_height*index)+0.01), Inches(pivot_column_width), Inches(standard_row_height))

def column_chart(slide, df, x, y, cx, cy):
    '''
    Creates stacked column chart
    '''
    ##Create chart object
    chart_data = ChartData()
    ##Define x-axis categories
    chart_data.categories = df['Year'].drop_duplicates()
    ##Create series for each grouping
    for series_name in df[df.columns[0]].unique():
        chart_data.add_series(series_name, df[df[df.columns[0]] == series_name]['Sales'])
    ##Add stacked column chart to slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data).chart
    ##Apply defined colors
    series_colors = [{'r':0,'g':48,'b':78},{'r':214,'g':40,'b':40},{'r':247,'g':127,'b':0},{'r':252,'g':191,'b':73}]
    ##Loop through each series and apply formatting
    for index, series_name in enumerate(df[df.columns[0]].unique()):
        ##Access series object
        series = chart.series[index]
        ##Apply fill color to series for legend
        _shape_fill(series.format, fill_color=series_colors[index])
        ##Loop through each point in the series and apply formatting
        for i, point in enumerate(series.points):
            ##Apply fill color to each point on the chart
            _shape_fill(point.format, series_colors[index])
            ##Apply border color to each point
            point.format.line.color.rgb = RGBColor(series_colors[index]['r'],series_colors[index]['g'],series_colors[index]['b'])
    ##Set y-axis font size and format
    chart.value_axis.tick_labels.font.size = util.Pt(12)
    chart.value_axis.tick_labels.number_format = '$#,##0'
    ##Set x-axis font size
    chart.category_axis.tick_labels.font.size = util.Pt(12)
    ##Add chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = 'Yearly Sales'
    ##Add legend
    chart.has_legend = True
    ##Position legend to the right of the chart
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = util.Pt(12)

def sparkline(slide, data, x, y, cx, cy):
    '''
    Creates sparkline chart
    '''
    chart_data = ChartData()
    chart_data.categories = [str(i) for i in range(len(data))]
    chart_data.add_series('Data', data)
    ##Add line chart to slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data).chart
    ##Remove plot elements
    _plot_area_only(chart)
    ##Access series object
    series = chart.series[0]
    ##Set line width and color
    series.format.line.width = util.Pt(1)
    series.format.line.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
    ##Add markers to each point
    series.marker.style = XL_MARKER_STYLE.CIRCLE
    ##Set marker size
    for point in series.points:
        point.marker.size = 2
    ##Apply formatting to markers
    _shape_fill(series.marker.format, fill_color=primary_color)
    _shape_fill(series.marker.format.line, fill_color=primary_color)

def sparkbar(slide, data, x, y, cx, cy):
    '''
    Creates sparkbar charts
    '''
    chart_data = ChartData()
    chart_data.categories = [str(i) for i in range(len(data))]
    chart_data.add_series('Data', data)
    ##Add column chart to slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    ##Set a minimum and maximum scale for consistency
    chart.value_axis.maximum_scale = max(data)
    chart.value_axis.minimum_scale = min(data)
    ##Remove plot elements
    _plot_area_only(chart)
    ##Create space between each column
    chart.plots[0].gap_width = 25
    ##Access series object
    series = chart.series[0]
    ##Loop through each point and apply formatting
    for i, point in enumerate(series.points):
        _ = SubElement(point.format.element, 'c:invertIfNegative', val=str(0))
        fill = point.format.fill
        fill.solid()
        ##If positive, set to green
        if series.values[i] > 0:
            fill.fore_color.rgb = RGBColor(0,176,80)
        ##If negative, set to red
        elif series.values[i] < 0:
            fill.fore_color.rgb = RGBColor(235,70,81)
        ##If zero, set to black
        else:
            fill.fore_color.rgb = RGBColor(0,0,0)
    
def seaborn_boxplot(slide, df, x, y, cx, cy):
    '''
    Creates boxplot with seaborn, saves the chart as an image file in-memory, then inserts the image onto the slide as a picture
    '''
    ##Create an in-memory binary stream that is stored in memory
    image_bytes = BytesIO()
    ##Create boxplot in seaborn
    sns.boxplot(
        data=df, y='Sales', x='Region', hue='Region',
        palette=['#00303E','#D62828','#F77F00','#FCBF49'],
        saturation=1
    )
    plt.ylabel('Sales')
    plt.xlabel('Region')
    plt.rcParams['figure.dpi'] = 600
    ##Save the boxplot as a png image
    plt.savefig(image_bytes, format='png', bbox_inches='tight')
    ##Access the binary file
    image_bytes.seek(0)
    ##Add binary file to slide
    slide.shapes.add_picture(image_bytes, x, y, cx, cy)

def linedivider(slide, x, y, cx, cy):
    '''
    Creates connector lines
    '''
    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x, y, cx, cy)
    _shape_color(shape, primary_color)
    shape.line.width = util.Pt(2.25)

#%%Randomly generated data
data = {
    'Region': list(itertools.chain.from_iterable([[item]*11 for item in ['Northeast','South','Midwest','West']])),
    'Year': [2015,2016,2017,2018,2019,2020,2021,2022,2023,2024,2025]*4,
    'Sales': list(np.random.uniform(low=10000, high=30000, size=44)),
}
df = pd.DataFrame(data)
##Pivot data to from long to wide for the table example
df_table = df.pivot(index=['Region'], columns='Year', values='Sales').reset_index(inplace=False)
##Add a totals row at the bottom of the dataframe
totals = df_table.sum(numeric_only=True)
totals['Region'] = 'Total'
df_table.loc['Total'] = totals
df_table = df_table.reset_index(inplace=False, drop=True)
##Calculate year-over-year values for sparkbar chart
df_table_qoq = df_table.drop(df_table.columns[1], axis=1).copy()
for col in df_table.select_dtypes(include='number').columns.to_list()[1:]:
    df_table_qoq[col] = ((df_table[col] - df_table[df_table.columns[df_table.columns.get_loc(col)-1]]) / df_table[df_table.columns[df_table.columns.get_loc(col)-1]])*100

#%%Presentation

##Create a new presentation object. You could also upload an existing template file by calling the file path: prs = Presentation('./template_file.pptx')
prs = Presentation()
##Change slide dimensions from the default
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
'''
Title slide:
Create a slide with both a main title and sub-title textbox
'''
##Create a new slide object. Select a slide layout index from list (base 0). 6 is for the 'Blank' slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
##Creates a new textbox object. The left and top kwargs represent the object's horizontal and vertical position.
textbox_title = slide.shapes.add_textbox(left=Inches(1), top=Inches(3.60), width=Inches(14), height=Inches(1))
##Inserts text into the textbox object
shape_text(textbox_title, ['Presentation Demo'], font_size=40, horizontal_alignment='center', vertical_alignment='middle')
##Sub-title textbox
textbox_subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.75), Inches(14), Inches(1))
shape_text(textbox_subtitle, ['python-pptx'], font_size=24, horizontal_alignment='center', vertical_alignment='middle')
'''
Text examples:
Provide examples for writing and formatting texts
'''
##Create a new slide object
slide = prs.slides.add_slide(prs.slide_layouts[6])
##Slide header
textbox = slide.shapes.add_textbox(left=Inches(0.05), top=Inches(0.10), width=Inches(9.80), height=Inches(0.40))
shape_text(textbox, text=['Text Example'], font_size=24, vertical_alignment='middle')
##Text frame
textbox = slide.shapes.add_textbox(left=Inches(0.93), top=Inches(0.86), width=Inches(4.34), height=Inches(1.00))
shape_text(textbox, text=['Textbox example. Each textbox has a single text frame and is made up of multiple paragraphs.'], font_size=14)
##Paragraphs
textbox = slide.shapes.add_textbox(left=Inches(0.93), top=Inches(2.10), width=Inches(4.34), height=Inches(1.92))
tf = textbox.text_frame
if tf.paragraphs[0]: tf.paragraphs[0]._element.getparent().remove(tf.paragraphs[0]._element)
for line in ['A paragraph is any block of text broken up by a new line character.','This is paragraph 2.','This is paragraph 3.','A text frame contains at least one initial paragraph.']:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = util.Pt(14)
    p.font.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
##Runs
textbox = slide.shapes.add_textbox(left=Inches(0.93), top=Inches(4.25), width=Inches(4.34), height=Inches(1.31))
shape_text(textbox, text=['A paragraph is made up of multiple runs. A run is a segment of text with shared styling.'], font_size=14)
##Horizontal alignment
textbox = slide.shapes.add_textbox(left=Inches(9.73), top=Inches(0.86), width=Inches(3.1), height=Inches(0.71))
shape_text(textbox, text=['Horizontal alignment is done at the paragraph level.'], horizontal_alignment='right', font_size=14)
_shape_color(textbox, line_color=primary_color)
##Vertical alignment
textbox = slide.shapes.add_textbox(left=Inches(9.73), top=Inches(1.76), width=Inches(4.34), height=Inches(1.30))
shape_text(textbox, text=['Vertical alignment is done at the text frame level.'], vertical_alignment='bottom', font_size=14)
_shape_color(textbox, line_color=primary_color)
##Formatting
textbox = slide.shapes.add_textbox(left=Inches(9.73), top=Inches(3.36), width=Inches(4.34), height=Inches(0.40))
shape_text(textbox, text=['Text formatting is done at the run level:'], font_size=14)
##Font color
textbox = slide.shapes.add_textbox(left=Inches(9.73), top=Inches(3.80), width=Inches(4.34), height=Inches(2.02))
shape_text(textbox, text=['Text color can be any RGB value. '], font_color={'r':255,'g':0,'b':0}, font_size=14)
tf = textbox.text_frame
p = tf.add_paragraph()
run = p.add_run()
run.text = 'Text emphasis such as '
run.font.size = util.Pt(14)
run.font.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
##Bold
run = p.add_run()
run.text = 'bold, '
run.font.size = util.Pt(14)
run.font.bold = True
run.font.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
##Italics
run = p.add_run()
run.text = 'italics, '
run.font.size = util.Pt(14)
run.font.italics = True
run.font.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
##Underline
run = p.add_run()
run.text = 'underline. '
run.font.size = util.Pt(14)
run.font.underline = True
run.font.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
##Font size
run = p.add_run()
run.text = 'Font size can be any point value. '
run.font.size = util.Pt(30)
run.font.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
##Hyperlink
run = p.add_run()
run.text = 'Text can be hyperlinked to a URL.'
run.font.size = util.Pt(14)
run.font.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
run.hyperlink.address = 'https://python-pptx.readthedocs.io/en/latest/'
'''
Shape Example:
Provide example of adding and formatting shape objects
'''
##Create a new slide object
slide = prs.slides.add_slide(prs.slide_layouts[6])
##Slide header
textbox = slide.shapes.add_textbox(left=Inches(0.05), top=Inches(0.10), width=Inches(9.80), height=Inches(0.40))
shape_text(textbox, text=['Shapes Example'], font_size=24, vertical_alignment='middle')
##Access the slide's shapes object
shapes = slide.shapes
##Create new shapes
autoshape(shapes, MSO_SHAPE.CLOUD, left=Inches(3.00), top=Inches(1.50), width=Inches(3.00), height=Inches(2.00), fill_color={'r':127,'g':127,'b':127}, line_width=util.Pt(3))
autoshape(shapes, MSO_SHAPE.LIGHTNING_BOLT, left=Inches(5.10), top=Inches(3.30), width=Inches(1.20), height=Inches(1.70), fill_color={'r':255,'g':255,'b':0})
##Rotate shape in degrees
autoshape(shapes, MSO_SHAPE.LIGHTNING_BOLT, left=Inches(6.30), top=Inches(1.90), width=Inches(1.20), height=Inches(1.70), fill_color={'r':255,'g':255,'b':0}, rotate=315)
##Remove default shadow effect
autoshape(shapes, MSO_SHAPE.LIGHTNING_BOLT, left=Inches(3.00), top=Inches(3.50), width=Inches(1.20), height=Inches(1.70), fill_color={'r':255,'g':255,'b':0}, rotate=55, shadow=False)
##Add text to shape
shape = autoshape(shapes, MSO_SHAPE.STAR_5_POINT, left=Inches(9.00), top=Inches(1.00), width=Inches(1.50), height=Inches(1.50))
shape_text(shape, text=['5'], font_size=32, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
shape = autoshape(shapes, MSO_SHAPE.STAR_6_POINT, left=Inches(9.00), top=Inches(3.00), width=Inches(1.50), height=Inches(1.50))
shape_text(shape, text=['6'], font_size=32, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
shape = autoshape(shapes, MSO_SHAPE.STAR_7_POINT, left=Inches(9.00), top=Inches(5.00), width=Inches(1.50), height=Inches(1.50))
shape_text(shape, text=['7'], font_size=32, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
shape = autoshape(shapes, MSO_SHAPE.STAR_8_POINT, left=Inches(9.00), top=Inches(7.00), width=Inches(1.50), height=Inches(1.50))
shape_text(shape, text=['8'], font_size=32, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
##Annotations
textbox = slide.shapes.add_textbox(left=Inches(0.40), top=Inches(0.95), width=Inches(4.03), height=Inches(0.81))
shape_text(textbox, text=['Shape fill/outline color along with outline weight. This cloud is colored dark gray with a white border that is 3px thick.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(6.71), top=Inches(0.81), width=Inches(1.64), height=Inches(1.52))
shape_text(textbox, text=['Shapes can be rotated by any degree. This lightning bolt is rotated 315 degrees.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(1.35), top=Inches(5.40), width=Inches(4.03), height=Inches(0.81))
shape_text(textbox, text=['Most shapes are created with a shadow effect by default. This lightning bolt has its shadow effect removed.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(10.93), top=Inches(1.44), width=Inches(3.75), height=Inches(1.04))
shape_text(textbox, text=['Most shape objects have an existing text fame that can be used to insert text. Each star below has the number of points written in the center of the shape.'], font_size=14)
'''
Chart Example:
Provide example of adding a PowerPoint chart
'''
##Create a new slide object
slide = prs.slides.add_slide(prs.slide_layouts[6])
##Slide header
textbox = slide.shapes.add_textbox(left=Inches(0.05), top=Inches(0.10), width=Inches(9.80), height=Inches(0.40))
shape_text(textbox, text=['Charts Example'], font_size=24, vertical_alignment='middle')
##Stacked column chart
column_chart(slide, df, x=Inches(0.30), y=Inches(2.00), cx=Inches(14.50), cy=Inches(4.30))
##Annotations
textbox = slide.shapes.add_textbox(left=Inches(0.51), top=Inches(0.95), width=Inches(13.56), height=Inches(0.34))
shape_text(textbox, text=['Most of the built-in PowerPoint charts are accessible in python. This is an example of a stacked column chart with randomly generated data. The plot contains 4 series, one for each region.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(0.43), top=Inches(6.37), width=Inches(6.51), height=Inches(0.57))
shape_text(textbox, text=['Axes labels can also be formatted, such as formatting the y-axis as $ amounts, or rotating the x-axis labels by 45/90 degrees for easier readability.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(13.73), top=Inches(6.43), width=Inches(2.15), height=Inches(0.57))
shape_text(textbox, text=['Chart legends can also be added/positioned.'], font_size=14)
'''
Table and Sparkline Example:
Provide example of adding a table and sparkline
'''
##Create a new slide object
slide = prs.slides.add_slide(prs.slide_layouts[6])
##Slide header
textbox = slide.shapes.add_textbox(left=Inches(0.05), top=Inches(0.10), width=Inches(9.80), height=Inches(0.40))
shape_text(textbox, text=['Tables and Sparklines Example'], font_size=24, vertical_alignment='middle')
##Main table
table(slide, df_table.iloc[:, [0]+list(range(-10,0))], title='Yearly Sales ($)', horizontal_pos=0.30, vertical_pos=2.00, font_size=14)
##Empty table where sparklines will be overlaid
trend_table(slide, df_table.iloc[:, list(range(-10,0))], first_line_title='10 Year', second_line_title='Trend', horizontal_pos=11.70, vertical_pos=2.00, chart_type='line', font_size=14)
trend_table(slide, df_table_qoq.iloc[:, list(range(-10,0))], first_line_title='10 Year', second_line_title='YoY Growth (%)', horizontal_pos=13.10, vertical_pos=2.00, chart_type='bar', font_size=14)
##Annotations
textbox = slide.shapes.add_textbox(left=Inches(0.30), top=Inches(0.82), width=Inches(8.18), height=Inches(0.57))
shape_text(textbox, text=["A table is created by setting the number of rows x columns. Each cell is like its own textbox and can have individual formatting. Cells can be accessed by looping through the table's index matrix."], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(0.70), top=Inches(1.60), width=Inches(6.66), height=Inches(0.34))
shape_text(textbox, text=['Cells can be merged together. The first row in the table is merged as a single cell.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(9.85), top=Inches(0.91), width=Inches(1.50), height=Inches(1.04))
shape_text(textbox, text=["A cell's background color can be any RGB value."], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(0.39), top=Inches(4.30), width=Inches(3.56), height=Inches(0.34))
shape_text(textbox, text=['Borders can be applied to individual cells.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(7.24), top=Inches(4.31), width=Inches(3.56), height=Inches(1.04))
shape_text(textbox, text=['Cell margins can be altered for better readability. All $ value cells have a left, top, and bottom cell margin of 0 inches with a small margin on the right.'], font_size=14)
textbox = slide.shapes.add_textbox(left=Inches(11.53), top=Inches(4.31), width=Inches(3.56), height=Inches(1.52))
shape_text(textbox, text=['Sparklines can be recreated by inserting a table with blank values and overlaying a small line chart to fit within the cell dimensions. All other plot elements are removed such as axes, labels, gridlines, and legends.'], font_size=14)
'''
Picture Example:
Provide example of adding a picture to a slide
'''
##Create a new slide object
slide = prs.slides.add_slide(prs.slide_layouts[6])
##Slide header
textbox = slide.shapes.add_textbox(left=Inches(0.05), top=Inches(0.10), width=Inches(9.80), height=Inches(0.40))
shape_text(textbox, text=['Pictures Example'], font_size=24, vertical_alignment='middle')
##Seaborn boxplot to save as image. If image save fails, provide alternative description
try:
    seaborn_boxplot(slide, df, x=Inches(0.30), y=Inches(2.00), cx=Inches(10.00), cy=Inches(5.00))
    textbox = slide.shapes.add_textbox(left=Inches(0.30), top=Inches(0.82), width=Inches(8.18), height=Inches(0.57))
    shape_text(textbox, text=['Any image file can be inserted into PowerPoint. For example, boxplots are not available in the python API, so the chart can be made in seaborn then saved as a PNG image and inserted onto the slide.'], font_size=14)
except:
    textbox = slide.shapes.add_textbox(left=Inches(0.30), top=Inches(0.82), width=Inches(8.18), height=Inches(0.57))
    shape_text(textbox, text=['Any image file can be inserted into PowerPoint. The image file was unable to be inserted into the file. Check the seaborn_boxplot() function.'], font_size=14)
'''
Conclusion
'''
##Create a new slide object
slide = prs.slides.add_slide(prs.slide_layouts[6])
##Slide header
textbox = slide.shapes.add_textbox(left=Inches(0.05), top=Inches(0.10), width=Inches(9.80), height=Inches(0.40))
shape_text(textbox, text=['Conclusion'], font_size=24, vertical_alignment='middle')
##Main paragraph text
conclusion_text = f'''
python-pptx is a versatile library that allows you to manipulate PowerPoint files. It is useful for analyses that require many different cuts of the data or have a refresh cycle (i.e., a slide deck is updated at the end of every month/quarter). PowerPoint files can be created using the default settings, or an existing file can be read and manipulated. Uploading an existing file can be useful if a company has a set template already defined.
This allows python to do a lot of the heavy lifting by calculating metrics and creating visuals quickly. Then the user can open the file and finish the presentation with any additional comments or formatting.
'''
textbox = slide.shapes.add_textbox(left=Inches(0.30), top=Inches(1.31), width=Inches(14.33), height=Inches(2.00))
shape_text(textbox, text=[conclusion_text], font_size=18)

##Save presentation in current folder
prs.save('./python_pptx_demo.pptx')

#%%
