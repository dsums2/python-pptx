#%%Libraries
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from io import BytesIO
from pptx import Presentation, util
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.chart.shared import CT_Layout
from pptx.chart.data import ChartData

#%%Config
primary_color = {'r':7,'g':24,'b':45}               # Navy Blue
secondary_color = {'r':241,'g':221,'b':157}         # Gold
white_color = {'r':255,'g':255,'b':255}             # White
positive_color = {'r':204,'g':255,'b':204}          # Green
positive_accent_color = {'r':0,'g':108,'b':49}      # Dark Green
negative_color = {'r':255,'g':209,'b':209}          # Red
negative_accent_color = {'r':155,'g':17,'b':27}     # Dark Red
standard_column_width = 0.70
pivot_column_width = 1.30
standard_row_height = 0.20

slide_order = [
    {'type':'transition', 'title':'Superstore Analysis', 'sub_title':'a python-pptx use case'},
    {'type':'pivot_summary','pivot':'Segment','title':'Segment'},
    {'type':'pivot_summary','pivot':'Region','title':'Region'},
    {'type':'pivot_summary','pivot':'Category','title':'Category'},
    {'type':'top_customers', 'title':'Top Performing Customers'},
    {'type':'ship_mode_analysis', 'title':'Ship Mode Comparison'},
    {'type':'transition','title':'Appendix'},
    {'type':'pivot_dbl_click','pivot':'State','title':'Region: East','filters':f'''Region == "East"'''},
    {'type':'pivot_dbl_click','pivot':'State','title':'Region: West','filters':f'''Region == "West"'''},
    {'type':'pivot_dbl_click','pivot':'State','title':'Region: South','filters':f'''Region == "South"'''},
    {'type':'pivot_dbl_click','pivot':'State','title':'Region: Central','filters':f'''Region == "Central"'''},
    {'type':'pivot_dbl_click','pivot':'Sub-Category','title':'Category: Furniture','filters':f'''Sub-Category == "Furniture"'''},
    {'type':'pivot_dbl_click','pivot':'Sub-Category','title':'Category: Technology','filters':f'''Sub-Category == "Technology"'''},
    {'type':'pivot_dbl_click','pivot':'Sub-Category','title':'Category: Office Supplies','filters':f'''Sub-Category == "Office Supplies"'''},
]

#%%Functions

def SubElement(parent, tagname, **kwargs):
    '''
    Function to define an element (common in XML syntax)
    '''
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _horizontal_alignment(shape, horizontal_alignment:str):
    '''
    Text align a paragraph horizontally.
    horizontal_alignment: alignment string; default is typically 'left' (dependent upon default settings)
    '''
    match horizontal_alignment:
        case 'left': shape.alignment = PP_ALIGN.LEFT
        case 'center': shape.alignment = PP_ALIGN.CENTER
        case 'right': shape.alignment = PP_ALIGN.RIGHT

def _vertical_alignment(shape, vertical_alignment:str):
    '''
    Text align a textframe vertically.
    vertical_alignment: alignment string; default is typically 'top' (dependent upon default settings)
    '''
    match vertical_alignment:
        case 'top': shape.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        case 'middle': shape.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        case 'bottom': shape.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM

def _shape_fill(shape, fill_color:dict):
    '''
    Fill shape with a solid color.
    fill_color: dictionary containing desired 'r', 'g', and 'b' color values
    '''
    ##Access shape's fill property
    fill = shape.fill
    ##Set fill pattern to solid
    fill.solid()
    ##Set shape's fill color as RGB value
    fill.fore_color.rgb = RGBColor(fill_color['r'],fill_color['g'],fill_color['b'])

def _shape_color(shape, line_color:dict):
    '''
    Fill shape outline with a solid color.
    line_color: dictionary containing desired 'r', 'g', and 'b' color values
    '''
    ##Access shape's line property
    line = shape.line
    ##Set shape's line color as RGB value
    line.color.rgb = RGBColor(line_color['r'],line_color['g'],line_color['b'])

def _paragraphs(shape, text:list, font_size:int, font_color:dict, bold:bool, horizontal_alignment:str, vertical_alignment:str, bulleted:bool, run_bold:list):
    '''
    Creates paragraph object
    '''
    ##Access shape's text frame
    tf = shape.text_frame
    ##A text frame has an initial paragraph object; we want to delete that and add our own
    if tf.paragraphs[0]: tf.paragraphs[0]._element.getparent().remove(tf.paragraphs[0]._element)
    ##Ensure text will maintain set width
    tf.word_wrap = True
    ##Loop through provided list of text; each item will be treated as its own paragraph
    for line in text:
        ##Add paragraph to shape's text frame
        p = tf.add_paragraph()
        ##If the item is also a list, break sub-items into separate runs to apply different text formatting
        if isinstance(line, list):
            for index, sub_line in enumerate(line):
                ##Add run to shape's paragraph
                run = p.add_run()
                ##Insert text as string
                run.text = sub_line
                run.font.size = util.Pt(font_size)
                run.font.bold = run_bold[index]
                run.font.color.rgb = RGBColor(font_color['r'],font_color['g'],font_color['b'])
        else:
            run = p.add_run()
            run.text = line
            run.font.size = util.Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(font_color['r'],font_color['g'],font_color['b'])
        ##If bulleted=True, make each paragraph in text frame a bullet point
        if bulleted:
            ##Access paragraph properties
            pPr = p._p.get_or_add_pPr()
            ##Set indentation between text frame and bullet; number is in EMU
            pPr.set('marL','0')
            ##Set indentation between bullet and paragraph; number is in EMU
            pPr.set('indent','171450')
            ##Define bullet properties, including the character symbol
            _ = SubElement(parent=pPr, tagname='a:buChar', char='●')
        ##Set text alignment
        _horizontal_alignment(p, horizontal_alignment)
        _vertical_alignment(tf, vertical_alignment)

def _cell_margins(shape, margin_left:float=0, margin_right:float=0, margin_top:float=0, margin_bottom:float=0, vertical_alignment:str='middle'):
    '''
    Set table cell margins.
    margin_left, margin_right, margin_top, margin_bottom: inches of padding within a table cell
    vertical_alignment: alignment string; default is typically 'top' (dependent upon default settings)
    '''
    shape.margin_left = Inches(margin_left)
    shape.margin_right = Inches(margin_right)
    shape.margin_top = Inches(margin_top)
    shape.margin_bottom = Inches(margin_bottom)
    _vertical_alignment(shape, vertical_alignment)

def _cell_border(cell, borders:list, line_type:str='solid', border_color:str="07182D", border_width:str='12700'):
    '''
    Apply table cell borders.
    borders: desired cell borders being accessed ('a:lnL'= left, 'a:lnR'= right, 'a:lnT'= top, 'a:lnB'= bottom); 'a:lnB' is the only side that works consistently
    line_type: border style ('solid', 'dash', 'dot', 'lgDash', 'dashDot', 'lgDashDot', 'sysDash', 'sysDot', 'sysDashDot', 'sysDashDotDot')
    border_color: border color; any hex code value
    border_width: string of English Metric Units (EMU), a common size standardization; 1px = 12700 EMUs, 1in = 914400 EMUs
    '''
    fill_found = False
    fill = cell.fill
    ##Check if cell already has a fill color; to properly apply a cell border, the background must be transparent
    if cell.fill.type == MSO_FILL.SOLID:
        ##Store fill color to reapply later
        fill_color = fill.fore_color.rgb
        fill_found = True
        ##Set cell fill to transparent
        fill.background()
    ##Access cell attributes
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    ##Clean up existing cell formatting by removing default cell attributes
    for lines in borders:
        tag = lines.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag):
                tcPr.remove(e)
        ##Apply desired cell borders
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        ##Apply border color and style
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val=line_type)
        ##Additional formatting options
        # round_ = SubElement(ln, 'a:round')
        # headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        # tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
    ##Reapply fill color if any was present
    if fill_found:
        fill.solid()
        cell.fill.fore_color.rgb = fill_color
    return cell

def _table_styles(shape):
    '''
    Apply specific table styling.
    Desired style is 'No Style, No Grid'
    Additional style GUIDs can be found at: https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372
    '''
    table_style = shape._element.graphic.graphicData.tbl
    table_style[0][-1].text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'

def _table_column_sizes(table, widths):
    '''
    Set each column width to the corresponding list of sizes
    '''
    table_columns = table.columns
    for index, width in enumerate(widths):
        table_columns[index].width = Inches(width)

def _table_row_sizes(table, heights):
    '''
    Set each row height to the corresponding list of sizes
    '''
    table_rows = table.rows
    for index, height in enumerate(heights):
        table_rows[index].height = Inches(height)

def _plot_area_only(chart):
    '''
    Removes all additional plot elements.
    Used when creating sparklines
    '''
    ##Remove y-axis gridlines
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    ##Remove y-axis title
    chart.value_axis.has_title = False
    ##Hide y-axis
    chart.value_axis.visible = False
    ##Set y-labels to a small font size and color same as slide background
    chart.value_axis.tick_labels.font.size = util.Pt(1)
    chart.value_axis.tick_labels.font.color.rgb = RGBColor(255,255,255)
    ##Remove x-axis gridlines
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    ##Remove x-axis title
    chart.category_axis.has_title = False
    ##Hide x-axis
    chart.category_axis.visible = False
    ##Set x-labels to a small font size and color same as slide background
    chart.category_axis.tick_labels.font.size = util.Pt(1)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor(255,255,255)
    ##Remove x-labels
    chart.category_axis.has_labels = False
    ##Remove chart title
    chart.has_title = False
    ##Remove legend
    chart.has_legend = False
    ##Remove any other plot element
    plot_area = chart.element.chart.plotArea
    plot_area_children = plot_area.getchildren()
    pac_layout = [x for x in plot_area_children if isinstance(x, CT_Layout)]
    if len(pac_layout) == 1:
        layout = pac_layout[0]
    else:
        layout = SubElement(plot_area, 'c:layout', val='inner')
    manual_layout = layout.get_or_add_manualLayout()
    manual_layout.remove_all(*['c:layoutTarget','c:xMode','c:yMode','c:x','c:y','c:w','c:h'])
    _ = SubElement(manual_layout, 'c:layoutTarget', val='inner')
    _ = SubElement(manual_layout, 'c:xMode', val='edge')
    _ = SubElement(manual_layout, 'c:yMode', val='edge')
    _ = SubElement(manual_layout, 'c:x', val=str(0))
    _ = SubElement(manual_layout, 'c:y', val=str(0))
    _ = SubElement(manual_layout, 'c:w', val=str(1))
    _ = SubElement(manual_layout, 'c:h', val=str(1))

def autoshape(shapes, mso_shape, left:float, top:float, width:float, height:float, fill_color:dict=primary_color, line_color:dict=white_color, no_border:bool=False, line_width:float=util.Pt(1), rotate:float=0, shadow:bool=True):
    '''
    Create an autoshape object.
    mso_shape: desired autoshape option
    no_border: if true, the outline will be transparent
    line_width: border thickness
    rotate: number of degrees to rotate the autoshape
    shadow: whether or not to include the default shadow effect
    '''
    shape = shapes.add_shape(mso_shape, left, top, width, height)
    _shape_fill(shape, fill_color=fill_color)
    _shape_color(shape, line_color=line_color)
    if no_border: shape.line.fill.background()
    shape.shadow.inherit = shadow
    shape.line.width = line_width
    shape.rotation = rotate
    return shape

def shape_text(shape, text:list, font_size:int, bold:bool=False, font_color:dict=primary_color, fill_color=False, border_color=False, horizontal_alignment:str='left', vertical_alignment:str='top', bulleted:bool=False, run_bold:list=[False,False]):
    '''
    Add text to a shape object.
    bulleted: whether or not to make each line of text in list as a bulleted item
    run_bold: bold two different runs in a bulleted line; useful for bolding a bulleted topic, and unbold its definition
    '''
    _paragraphs(shape, text, font_size, font_color, bold, horizontal_alignment, vertical_alignment, bulleted, run_bold)
    if border_color: _shape_color(shape, border_color)
    if fill_color: _shape_fill(shape, fill_color)

def value_table(df, title, horizontal_pos, vertical_pos):
    '''
    Creates table object for sales $
    '''
    ##Create table with (number of rows, number of columns, left position, top position, width, height)
    shape = slide.shapes.add_table(len(df)+2, len(df.columns), Inches(horizontal_pos), Inches(vertical_pos), Inches((standard_column_width*(len(df.columns)-1))+pivot_column_width), Inches(standard_row_height*(len(df)+2)))
    table = shape.table
    ##Apply table styling
    _table_styles(shape)
    ##Set column widths. First column is wider and the rest are standard column width
    _table_column_sizes(table, [pivot_column_width]+[standard_column_width]*(len(df.columns)-1))
    ##Access first cell in table (base 0)
    cell = table.cell(0,0)
    ##Access last cell in the first row
    other_cell = table.cell(0,len(df.columns)-1)
    ##Merge all cells between them into a single cell
    cell.merge(other_cell)
    ##Set cell margins
    _cell_margins(cell)
    ##Cell color background
    _shape_fill(cell, fill_color=primary_color)
    ##Insert text into cell
    shape_text(cell, [title], 10, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
    ##Loop through each cell in the dataframe and insert into the table
    for i in range(1,len(df)+2):
        for j in range(0,len(df.columns)):
            cell = table.cell(i,j)
            _cell_margins(cell)
            ##If the cell is in the second row, insert column header
            if i == 1:
                _shape_fill(cell, fill_color=primary_color)
                shape_text(cell, [str(df.columns[j]).replace('_', ' ')], 10, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
            ##If the dataframe cell is a number, apply certain formatting
            elif isinstance(df.iloc[i-2,j], (int, float)):
                cell.margin_right = Inches(0.05)
                shape_text(cell, [f'$ {df.iloc[i-2,j]:,.0f}'], 10, bold=i-1==len(df), horizontal_alignment='right', vertical_alignment='middle')
                ##If cell is in the last or second to last row, apply a bottom cell border
                if (i-1 == len(df)) or (i-1 == len(df)-1): cell = _cell_border(cell, ['a:lnB'])
            ##If the dataframe cell is a string, apply generic formatting
            else:
                cell.margin_left = Inches(0.05)
                shape_text(cell, [str(df.iloc[i-2,j])], 10, bold=i-1==len(df), horizontal_alignment='left', vertical_alignment='middle')
                if (i-1 == len(df)) or (i-1 == len(df)-1): cell = _cell_border(cell, ['a:lnB'])
    ##Set row heights. Must be set after text is inserted into all cells. When text is initially inserted into a cell it has the default font size; if the desired row height is smaller than default font size, the cell will not get smaller.
    _table_row_sizes(table, [standard_row_height]*(len(df)+2))

def qoq_table(df, title, horizontal_pos, vertical_pos):
    '''
    Creates table object for sales growth
    '''
    shape = slide.shapes.add_table(len(df)+2, len(df.columns), Inches(horizontal_pos), Inches(vertical_pos), Inches((standard_column_width*(len(df.columns)-1))+pivot_column_width), Inches(standard_row_height*(len(df)+2)))
    table = shape.table
    _table_styles(shape)
    _table_column_sizes(table, [pivot_column_width]+[standard_column_width]*(len(df.columns)-1))

    cell = table.cell(0,0)
    other_cell = table.cell(0,len(df.columns)-1)
    cell.merge(other_cell)
    _cell_margins(cell)
    _shape_fill(cell, fill_color=primary_color)
    shape_text(cell, [title], 10, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')

    for i in range(1,len(df)+2):
        for j in range(0,len(df.columns)):
            cell = table.cell(i,j)
            _cell_margins(cell)
            if i == 1:
                text = str(df.columns[j]).replace('_', ' ')
                horizontal_alignment = 'center'
                bold = True
                fill_color = primary_color
                font_color = white_color
            elif isinstance(df.iloc[i-2,j], (int, float)):
                text = f'{df.iloc[i-2,j]:,.1f}%'
                horizontal_alignment = 'right'
                bold = i-1==len(df)
                cell.margin_right = Inches(0.05)
                ##If dataframe cell is negative, color red
                if round(df.iloc[i-2,j],2) < 0:
                    fill_color = negative_color
                    font_color = negative_accent_color
                ##If dataframe cell is positive, color green
                elif round(df.iloc[i-2,j],2) > 0:
                    fill_color = positive_color
                    font_color = positive_accent_color
                ##If dataframe cell is 0, apply primary color
                else:
                    fill_color = white_color
                    font_color = primary_color
            else:
                cell.margin_left = Inches(0.05)
                text = str(df.iloc[i-2,j])
                horizontal_alignment = 'left'
                bold = i-1==len(df)
                fill_color = white_color
                font_color = primary_color
            _shape_fill(cell, fill_color=fill_color)
            shape_text(cell, [text], 10, bold=bold, font_color=font_color, horizontal_alignment=horizontal_alignment, vertical_alignment='middle')
            if (i-1 == len(df)) or (i-1 == len(df)-1): cell = _cell_border(cell, ['a:lnB'])
    _table_row_sizes(table, [standard_row_height]*(len(df)+2))

def trend_table(df, first_line_title, second_line_title, horizontal_pos, vertical_pos, chart_type, inverse:bool=False):
    '''
    Creates placeholder table where sparklines will be overlaid
    '''
    shape = slide.shapes.add_table(len(df)+2, 1, Inches(horizontal_pos), Inches(vertical_pos), Inches(pivot_column_width), Inches(standard_row_height*(len(df)+2)))
    table = shape.table
    _table_styles(shape)
    _table_column_sizes(table, [pivot_column_width])
    ##Table headers
    for index, title in enumerate([first_line_title, second_line_title]):
        cell = table.cell(index, 0)
        _cell_margins(cell)
        _shape_fill(cell, fill_color=primary_color)
        shape_text(cell, [title], 10, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
    ##To maintain desired formatting, text must be in every cell; '0' is inserted into each cell and colored the same as the slide background to appear hidden.
    for i in range(2, len(df)+2):
        cell = table.cell(i,0)
        _cell_margins(cell)
        shape_text(cell, ['0'], 10, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
        if (i==len(df)) or (i==len(df)+1): cell = _cell_border(cell, ['a:lnB'])
    _table_row_sizes(table, [standard_row_height]*(len(df)+2))
    ##Overlay sparkline/sparkbar in the table cell
    for index, row in df.iterrows():
        if chart_type == 'line': sparkline(slide, [item for item in row.tolist()[-12:] if isinstance(item, (int, float))], Inches(horizontal_pos), Inches(vertical_pos+(standard_row_height*2)+(standard_row_height*index)), Inches(pivot_column_width), Inches(standard_row_height))
        if chart_type == 'bar': sparkbar(slide, [item for item in row.tolist()[-12:] if isinstance(item, (int, float))], Inches(horizontal_pos), Inches(vertical_pos+(standard_row_height*2)+(standard_row_height*index)+0.01), Inches(1.18), Inches(0.12))

def prop_table(df, title, horizontal_pos, vertical_pos):
    '''
    Creates table object for sales proportion by region
    '''
    shape = slide.shapes.add_table(len(df)+2, 3, Inches(horizontal_pos), Inches(vertical_pos), Inches(standard_column_width*3), Inches(standard_row_height*(len(df)+2)))
    table = shape.table
    _table_styles(shape)
    _table_column_sizes(table, [standard_column_width]*(3))

    cell = table.cell(0,0)
    other_cell = table.cell(0,2)
    cell.merge(other_cell)
    _cell_margins(cell)
    _shape_fill(cell, fill_color=primary_color)
    shape_text(cell, [title], 10, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')

    for i in range(1,len(df)+2):
        for j in range(0,3):
            cell = table.cell(i,j)
            _cell_margins(cell)
            if i == 1:
                cell.margin_right = Inches(0)
                _shape_fill(cell, fill_color=primary_color)
                shape_text(cell, [str(df.columns[j+1]).replace('_',' ')], 10, bold=True, font_color=white_color, horizontal_alignment='center', vertical_alignment='middle')
            elif isinstance(df.iloc[i-2,j+1], (int, float)):
                cell.margin_right = Inches(0.05)
                if i-1 == len(df):
                    shape_text(cell, [f'{df.iloc[i-2,j+1]:,.1f}%'], 10, bold=True, horizontal_alignment='right', vertical_alignment='middle')
                    cell = _cell_border(cell, ['a:lnB'])
                else:
                    shape_text(cell, [f'{df.iloc[i-2,j+1]:,.1f}%'], 10, bold=False, horizontal_alignment='right', vertical_alignment='middle')
            else:
                cell.margin_left = Inches(0.05)
                if i-1 == len(df):
                    shape_text(cell, [str(df.iloc[i-2,j+1])], 10, bold=True, vertical_alignment='middle')
                    cell = _cell_border(cell, ['a:lnB'])
                else:
                    shape_text(cell, [str(df.iloc[i-2,j+1])], 10, bold=False, vertical_alignment='middle')
            if i-1 == len(df)-1:
                cell = _cell_border(cell, ['a:lnB'])
    _table_row_sizes(table, [standard_row_height]*(len(df)+2))
    ##Overlay proportion bar in the table cell
    for index, row in df.iterrows():
        propbar(row[f'{str(int(df.columns[-1][:4])-2)}-{df.columns[-1][-2:]}']/100, index, horizontal_pos+(standard_column_width*0), vertical_pos+(2*standard_row_height)+0.02)
        propbar(row[f'{str(int(df.columns[-1][:4])-1)}-{df.columns[-1][-2:]}']/100, index, horizontal_pos+(standard_column_width*1), vertical_pos+(2*standard_row_height)+0.02)
        propbar(row[f'{str(int(df.columns[-1][:4]))}-{df.columns[-1][-2:]}']/100, index, horizontal_pos+(standard_column_width*2), vertical_pos+(2*standard_row_height)+0.02)

def monthly_column_chart(slide, df, x, y, cx, cy):
    '''
    Creates stacked column chart
    '''
    ##Create chart object
    chart_data = ChartData()
    ##Define x-axis categories
    chart_data.categories = df['Calendar_Month'].drop_duplicates()
    ##Create series for each grouping
    for series_name in df[df.columns[0]].unique():
        chart_data.add_series(series_name, df[df[df.columns[0]] == series_name]['Sales'])
    ##Add stacked column chart to slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data).chart
    ##Apply defined colors
    series_colors = [{'r':0,'g':48,'b':78},{'r':214,'g':40,'b':40},{'r':247,'g':127,'b':0},{'r':252,'g':191,'b':73}]
    ##Loop through each series and apply formatting
    for index, series_name in enumerate(df[df.columns[0]].unique()):
        ##Access series object
        series = chart.series[index]
        ##Apply fill color to series for legend
        _shape_fill(series.format, fill_color=series_colors[index])
        ##Loop through each point in the series and apply formatting
        for i, point in enumerate(series.points):
            ##Apply fill color to each point on the chart
            _shape_fill(point.format, series_colors[index])
            ##Apply border color to each point
            point.format.line.color.rgb = RGBColor(series_colors[index]['r'],series_colors[index]['g'],series_colors[index]['b'])
    ##Set y-axis font size and format
    chart.value_axis.tick_labels.font.size = util.Pt(12)
    chart.value_axis.tick_labels.number_format = '$#,##0'
    ##Set x-axis font size
    chart.category_axis.tick_labels.font.size = util.Pt(12)
    ##Add chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = 'Monthly Sales'
    ##Add legend
    chart.has_legend = True
    ##Position legend to the right of the chart
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = util.Pt(12)

def sparkline(slide, data, x, y, cx, cy):
    '''
    Creates sparkline chart
    '''
    chart_data = ChartData()
    chart_data.categories = [str(i) for i in range(len(data))]
    chart_data.add_series('Data', data)
    ##Add line chart to slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data).chart
    ##Remove plot elements
    _plot_area_only(chart)
    ##Access series object
    series = chart.series[0]
    ##Set line width and color
    series.format.line.width = util.Pt(1)
    series.format.line.color.rgb = RGBColor(primary_color['r'],primary_color['g'],primary_color['b'])
    ##Add markers to each point
    series.marker.style = XL_MARKER_STYLE.CIRCLE
    ##Set marker size
    for point in series.points:
        point.marker.size = 2
    ##Apply formatting to markers
    _shape_fill(series.marker.format, fill_color=primary_color)
    _shape_fill(series.marker.format.line, fill_color=primary_color)

def sparkbar(slide, data, x, y, cx, cy):
    '''
    Creates sparkbar charts
    '''
    chart_data = ChartData()
    chart_data.categories = [str(i) for i in range(len(data))]
    chart_data.add_series('Data', data)
    ##Add column chart to slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    ##Set a minimum and maximum scale for consistency
    chart.value_axis.maximum_scale = max(data)
    chart.value_axis.minimum_scale = min(data)
    ##Remove plot elements
    _plot_area_only(chart)
    ##Create space between each column
    chart.plots[0].gap_width = 25
    ##Access series object
    series = chart.series[0]
    ##Loop through each point and apply formatting
    for i, point in enumerate(series.points):
        _ = SubElement(point.format.element, 'c:invertIfNegative', val=str(0))
        fill = point.format.fill
        fill.solid()
        ##If positive, set to green
        if series.values[i] > 0:
            fill.fore_color.rgb = RGBColor(0,176,80)
        ##If negative, set to red
        elif series.values[i] < 0:
            fill.fore_color.rgb = RGBColor(235,70,81)
        ##If zero, set to black
        else:
            fill.fore_color.rgb = RGBColor(0,0,0)

def propbar(width, index, x, y):
    '''
    Proportion bar overlaid in the prop table
    '''
    shapes = slide.shapes
    ##Add rectangle shape
    shape = autoshape(shapes, MSO_SHAPE.RECTANGLE, left=Inches(x), top=Inches(y+(standard_row_height*index)), width=Inches(0.69*width), height=Inches(0.16), fill_color=secondary_color, no_border=True, shadow=False)
    ##Move rectangle behind table
    base_shape = shapes[0]._element
    base_shape.addprevious(shape._element)

def seaborn_boxplot(slide, df, x, y, cx, cy):
    '''
    Creates boxplot with seaborn, saves the chart as an image file in-memory, then inserts the image onto the slide as a picture
    '''
    ##Create an in-memory binary stream that is stored in memory
    image_bytes = BytesIO()
    ##Create boxplot in seaborn
    sns.boxplot(
        data=df, y='Delivery_Length', x='Ship_Mode', hue='Ship_Mode',
        palette=['#FCBF49','#00303E','#D62828','#F77F00'],
        saturation=1
    )
    plt.ylabel('Delivery Length (in days)')
    plt.xlabel('Ship Mode')
    plt.rcParams['figure.dpi'] = 600
    ##Save the boxplot as a png image
    plt.savefig(image_bytes, format='png', bbox_inches='tight')
    ##Access the binary file
    image_bytes.seek(0)
    ##Add binary file to slide
    slide.shapes.add_picture(image_bytes, x, y, cx, cy)

def linedivider(slide, x, y, cx, cy):
    '''
    Creates connector lines
    '''
    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x, y, cx, cy)
    _shape_color(shape, primary_color)
    shape.line.width = util.Pt(2.25)

#%%Data
##Read data from CSV
df = pd.read_csv('./superstore.csv').drop(columns=['Row ID'])
df.columns = df.columns.str.replace(' ', '_')
df['Order_Date'] = pd.to_datetime(df['Order_Date'], format='%d/%m/%Y')
df['Ship_Date'] = pd.to_datetime(df['Ship_Date'], format='%d/%m/%Y')
df['Delivery_Length'] = (df['Ship_Date']-df['Order_Date']).dt.days
df['Order_Month'] = df['Order_Date'].dt.month
df['Order_Quarter'] = np.ceil(df['Order_Month']/3).astype(int)
df['Order_Year'] = df['Order_Date'].dt.year
df['Calendar_Quarter'] = df['Order_Year'].astype(str) + '-Q' + df['Order_Quarter'].astype(str)
df['Calendar_Month'] = pd.to_datetime(df['Order_Month'], format='%m').dt.strftime('%b') + '-' + df['Order_Year'].astype(str)

#%%Presentation
##Create a new presentation object. You could also upload an existing template file by calling the file path: prs = Presentation('./template_file.pptx')
prs = Presentation()
##Change slide dimensions from the default
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
##Loop through each dictionary in slide order list
for current_slide in slide_order:
    ##Calculate necessary dataframes for slide
    match current_slide['type']:
        case x if x.startswith('pivot'):
            ## Dataframes
            try:
                df_hc = df.query(current_slide['filters'])
            except:
                df_hc = df.copy()
            df_hc = df_hc[[current_slide['pivot'],'Calendar_Quarter','Sales']].groupby([current_slide['pivot'],'Calendar_Quarter']).sum().reset_index(inplace=False).pivot(index=[current_slide['pivot']], columns='Calendar_Quarter', values='Sales').reset_index(inplace=False)
            df_hc = df_hc.fillna(0)

            ## Add total rows at bottom of dataframe
            totals = df_hc.sum(numeric_only=True)
            totals[current_slide['pivot']] = 'Total'
            df_hc.loc['Total'] = totals
            df_hc = df_hc.reset_index(inplace=False, drop=True)

            ## Calculate quarter-over-quarter growth for HC and 4 quarter rolling attrition for ATTR
            df_hc_qoq = df_hc.drop(df_hc.columns[1:5], axis=1).copy()
            for col in df_hc.select_dtypes(include='number').columns.to_list()[4:]:
                df_hc_qoq[col] = ((df_hc[col] - df_hc[df_hc.columns[df_hc.columns.get_loc(col)-4]]) / df_hc[df_hc.columns[df_hc.columns.get_loc(col)-4]])*100

            ## Calculate pivot column mix percentage
            df_hc_mix = df_hc[[current_slide['pivot']]+list(df_hc.columns[df_hc.columns.str.endswith(df_hc.columns[-1][-2:])])].copy()
            for col in df_hc_mix.select_dtypes(include='number').columns.to_list():
                df_hc_mix[col] = (df_hc_mix[col]/df_hc_mix[col].iloc[-1])*100
            df_hc_mix = df_hc_mix.drop(df_hc_mix.columns[1], axis=1)
            df_hc = df_hc.drop(df_hc.columns[1:5], axis=1)

            ## Clean up dataframes
            df_hc_qoq = df_hc_qoq.fillna(0)
            df_hc_mix = df_hc_mix.fillna(0)
            df_hc_qoq.replace([np.inf], 100, inplace=True)
            ##If pivot_summary, then include a monthly sales chart at the bottom
            if current_slide['type'] == 'pivot_summary': df_monthly = df[df['Order_Year'] != 2015][[current_slide['pivot'],'Order_Year','Order_Month','Calendar_Month','Sales']].groupby([current_slide['pivot'],'Order_Year','Order_Month','Calendar_Month']).sum().sort_values(by=[current_slide['pivot'],'Order_Year','Order_Month']).reset_index(inplace=False)
        case 'top_customers':
            ##Get top 3 customers
            df_hc = df[df['Order_Year'] != 2015][['Customer_Name','Customer_ID','Segment','Sales']].groupby(['Customer_Name','Customer_ID','Segment']).sum().sort_values(by=['Sales'], ascending=[False]).head(3).rename(columns={'Sales':'Total_Sales'}).reset_index(inplace=False)
            ##Merge customer information for top 3 customers
            df_hc = df_hc.merge(df[['Customer_Name','Order_Date']].groupby(['Customer_Name']).min().rename(columns={'Order_Date':'First_Order_Date'}), how='left', on='Customer_Name')
            ##Merge first month where the customer made an order
            df_hc['First_Calendar_Month'] = df_hc['First_Order_Date'].dt.strftime('%b-%Y')
            ##Merge number of orders made by each customer
            df_hc = df_hc.merge(df[df['Order_Year'] != 2015][['Customer_Name','Order_ID']].drop_duplicates().groupby(['Customer_Name']).count().rename(columns={'Order_ID':'Order_ID_Count'}), how='left', on='Customer_Name')
            ##Merge yearly sales average made by each customer
            df_hc = df_hc.merge(df[df['Order_Year'] != 2015][['Customer_Name','Order_Year','Sales']].groupby(['Customer_Name','Order_Year']).sum().reset_index(inplace=False)[['Customer_Name','Sales']].groupby(['Customer_Name']).mean().rename(columns={'Sales':'Sales_Yearly_Average'}), how='left', on='Customer_Name')
            ##Determine the average amount of months between orders for each customer
            temp = df[['Customer_Name','Order_ID','Order_Date']].drop_duplicates().groupby(['Customer_Name','Order_ID']).min().reset_index(inplace=False).sort_values(by=['Customer_Name','Order_Date'])
            temp['date_diff'] = temp.groupby(['Customer_Name'])['Order_Date'].diff().dt.days
            df_hc = df_hc.merge(temp.groupby(['Customer_Name'])['date_diff'].mean().reset_index(inplace=False).rename(columns={'date_diff':'Average_Order_Length_Days'}), how='left', on='Customer_Name')
            df_hc['Average_Order_Length_Months'] = df_hc['Average_Order_Length_Days']/30
            ##Determine the most expensive products bought by each customer
            df_top_products = df[df['Order_Year'] != 2015][['Customer_Name','Product_Name','Sales']].groupby(['Customer_Name','Product_Name']).sum().reset_index(inplace=False).sort_values(by=['Customer_Name','Sales'], ascending=[True, False])
            df_top_products['Sales'] = df_top_products.apply(lambda row: f'  (${row['Sales']:,.0f})', axis=1)
            ##Concatenate product name and sales cost as a list
            df_top_products['Product_Sales_Concat'] = df_top_products[['Product_Name','Sales']].values.tolist()
            df_top_products = df_top_products[['Customer_Name','Product_Sales_Concat']].groupby(['Customer_Name']).agg(list).reset_index(inplace=False)
            ##Limit list to the top 4 products
            df_top_products['Top_Products_Purchased'] = df_top_products['Product_Sales_Concat'].apply(lambda row: row[:4])
            df_hc = df_hc.merge(df_top_products, how='left', on='Customer_Name')
        case 'ship_mode_analysis':
            df_ship_mode = df[df['Order_Year'] != 2015][['Order_ID','Ship_Mode','Delivery_Length','Product_Name']].groupby(['Order_ID','Ship_Mode','Delivery_Length']).count().reset_index(inplace=False).rename(columns={'Product_Name':'Order_Product_Count'})
            df_ship_mode_sorted = df_ship_mode.copy()
            df_ship_mode_sorted['Ship_Mode'] = pd.Categorical(df_ship_mode['Ship_Mode'], categories=['Same Day','First Class','Second Class','Standard Class'], ordered=True)
            df_ship_mode_sorted = df_ship_mode_sorted.sort_values(by=['Ship_Mode'])
    ##Create a new slide object. Select a slide layout index from list (base 0). 6 is for the 'Blank' slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    match current_slide['type']:
        case x if x.startswith('pivot'):
            ##Slide header
            textbox = slide.shapes.add_textbox(Inches(0.05), Inches(0.10), Inches(9.80), Inches(0.40))
            shape_text(textbox, [current_slide['title']], 24, vertical_alignment='middle')
            ##Count Table
            value_table(df_hc, title='Sales ($)', horizontal_pos=0.10, vertical_pos=0.64)
            ##QoQ Table
            qoq_table(df_hc_qoq, title='Sales YoY Growth (%)', horizontal_pos=0.10, vertical_pos=((standard_row_height*(len(df_hc)+3))+0.64))
            ##Trend Table
            trend_table(df_hc, first_line_title='12 Quarter', second_line_title='Trend', horizontal_pos=9.90, vertical_pos=0.64, chart_type='line')
            trend_table(df_hc_qoq, first_line_title='12 Quarter', second_line_title='YoY Growth (%)', horizontal_pos=9.90, vertical_pos=((standard_row_height*(len(df_hc)+3))+0.64), chart_type='bar')
            ##Proportion Table
            prop_table(df_hc_mix, title='Sales Proportion (%)', horizontal_pos=11.30, vertical_pos=0.64)
            ##Monthly Table
            if current_slide['type'] == 'pivot_summary': monthly_column_chart(slide, df_monthly, Inches(0.10), Inches((standard_row_height*(len(df_hc)+3))+(standard_row_height*(len(df_hc_qoq)+3))+0.64), Inches(14.50), Inches(4.30))
        case 'top_customers':
            ##Gold, Silver, bronze medal colors
            medal_color = [{'r':211,'g':175,'b':55},{'r':196,'g':196,'b':196},{'r':206,'g':137,'b':54}]
            ##Slide header
            textbox = slide.shapes.add_textbox(Inches(0.05), Inches(0.10), Inches(9.80), Inches(0.40))
            shape_text(textbox, [current_slide['title']], 24, vertical_alignment='middle')
            ##Loop through each top customer
            for index, customer in enumerate(df_hc['Customer_Name'].unique()):
                ##Ribbon
                shapes = slide.shapes
                autoshape(shapes, MSO_SHAPE.CHEVRON, Inches(-0.06), Inches(1.69+(2.48*index)), Inches(1.07), Inches(0.36), rotate=270)
                ##Medal
                shape = autoshape(shapes, MSO_SHAPE.STAR_7_POINT, Inches(0.12), Inches(0.99+(2.48*index)), Inches(0.70), Inches(0.70))
                shape_text(shape, [str(index+1)], 18, bold=True, font_color=white_color, fill_color=medal_color[index], border_color={'r':0,'g':0,'b':0}, horizontal_alignment='center', vertical_alignment='middle')
                ##Customer Name
                textbox = slide.shapes.add_textbox(Inches(0.91), Inches(1.14+(2.48*index)), Inches(4.02), Inches(0.4))
                shape_text(textbox, [customer], 18, vertical_alignment='middle')
                ##Total Sales ($)
                textbox = slide.shapes.add_textbox(Inches(0.91), Inches(1.60+(2.48*index)), Inches(1.67), Inches(0.64))
                shape_text(textbox, [f'${df_hc['Total_Sales'].iloc[index]:,.0f}'], 32, horizontal_alignment='center', vertical_alignment='middle')
                textbox = slide.shapes.add_textbox(Inches(0.91), Inches(2.14+(2.48*index)), Inches(1.67), Inches(0.30))
                shape_text(textbox, ['in the last 3 years'], 10, horizontal_alignment='center', vertical_alignment='middle')
                ##Customer Demographics
                textbox = slide.shapes.add_textbox(Inches(2.75), Inches(1.60+(2.48*index)), Inches(2.18), Inches(0.91))
                shape_text(textbox, [[f'Customer ID: ',f'{df_hc['Customer_ID'].iloc[index]}'],[f'Segment: ',f'{df_hc['Segment'].iloc[index]}'],[f'Customer Since: ',f'{df_hc['First_Calendar_Month'].iloc[index]}'],[f'Number of Orders: ',f'{df_hc['Order_ID_Count'].iloc[index]}']], 12, vertical_alignment='middle', bulleted=True, run_bold=[False,True])
                ##Divider
                linedivider(slide, Inches(5.13), Inches(1.60+(2.48*index)), Inches(0), Inches(1.10))
                ##Average Sales ($) per Year
                textbox = slide.shapes.add_textbox(Inches(5.30), Inches(1.60+(2.48*index)), Inches(1.67), Inches(0.64))
                shape_text(textbox, [f'${df_hc['Sales_Yearly_Average'].iloc[index]:,.0f}'], 32, horizontal_alignment='center', vertical_alignment='middle')
                textbox = slide.shapes.add_textbox(Inches(5.30), Inches(2.14+(2.48*index)), Inches(1.67), Inches(0.30))
                shape_text(textbox, ['on average, annually'], 10, horizontal_alignment='center', vertical_alignment='middle')
                ##Average Number of Months between Orders
                textbox = slide.shapes.add_textbox(Inches(6.94), Inches(1.60+(2.48*index)), Inches(1.67), Inches(0.64))
                shape_text(textbox, [f'{df_hc['Average_Order_Length_Months'].iloc[index]:,.2f}'], 32, horizontal_alignment='center', vertical_alignment='middle')
                textbox = slide.shapes.add_textbox(Inches(6.94), Inches(2.14+(2.48*index)), Inches(1.67), Inches(0.30))
                shape_text(textbox, ['months between orders'], 10, horizontal_alignment='center', vertical_alignment='middle')
                ##Divider
                linedivider(slide, Inches(8.78), Inches(1.60+(2.48*index)), Inches(0), Inches(1.10))
                ##Top Products
                textbox = slide.shapes.add_textbox(Inches(8.95), Inches(1.54+(2.48*index)), Inches(2.80), Inches(0.40))
                shape_text(textbox, ['Top Products Purchased'], 18, vertical_alignment='middle')
                textbox = slide.shapes.add_textbox(Inches(8.95), Inches(1.90+(2.48*index)), Inches(6.30), Inches(0.85))
                shape_text(textbox, df_hc['Top_Products_Purchased'].iloc[index], 11, vertical_alignment='middle', bulleted=True, run_bold=[False,True])
        case 'ship_mode_analysis':
            ##Slide header
            textbox = slide.shapes.add_textbox(Inches(0.05), Inches(0.10), Inches(9.80), Inches(0.40))
            shape_text(textbox, [current_slide['title']], 24, vertical_alignment='middle')
            ##Seaborn boxplot to save as image. If image save fails, provide alternative description
            try:
                seaborn_boxplot(slide, df_ship_mode_sorted, x=Inches(0.10), y=Inches(1.00), cx=Inches(9.00), cy=Inches(7.00))
            except:
                textbox = slide.shapes.add_textbox(left=Inches(0.30), top=Inches(1.00), width=Inches(8.18), height=Inches(0.57))
                shape_text(textbox, text=['The image file was unable to be inserted into the file.'], font_size=14)
        case 'transition':
            ##Slide header
            textbox = slide.shapes.add_textbox(Inches(1), Inches(3.60), Inches(14), Inches(1))
            shape_text(textbox, [current_slide['title']], 40, horizontal_alignment='center', vertical_alignment='middle')
            if 'sub_title' in current_slide:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(4.75), Inches(14), Inches(1))
                shape_text(textbox, [current_slide['sub_title']], 24, horizontal_alignment='center', vertical_alignment='middle')
##Save presentation in current folder
prs.save('./superstore_example.pptx')

#%%
